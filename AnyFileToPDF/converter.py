import os
import mimetypes
import chardet
from PIL import Image
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import threading
import logging
from pathlib import Path
import io
import html
import re

class PDFConverter:
    def __init__(self):
        self.cancel_flag = False
        self.logger = logging.getLogger(__name__)
        self.setup_fonts()
        self.setup_styles()
        self.setup_mime_types()
        
    def setup_fonts(self):
        """设置字体"""
        try:
            # 注册中文字体（如果系统中有的话）
            font_paths = [
                "C:/Windows/Fonts/msyh.ttf",  # 微软雅黑
                "C:/Windows/Fonts/simsun.ttc", # 宋体
                "C:/Windows/Fonts/simhei.ttf", # 黑体
                "C:/Windows/Fonts/STSONG.TTF", # 华文宋体
                "C:/Windows/Fonts/STKAITI.TTF", # 华文楷体
            ]
            
            for font_path in font_paths:
                if os.path.exists(font_path):
                    font_name = os.path.splitext(os.path.basename(font_path))[0]
                    try:
                        pdfmetrics.registerFont(TTFont(font_name, font_path))
                        self.default_font = font_name
                        self.logger.info(f"成功注册字体: {font_name}")
                        return
                    except Exception as e:
                        self.logger.warning(f"注册字体失败 {font_name}: {str(e)}")
                        continue
                        
            # 如果没有找到中文字体，使用默认的Helvetica
            self.default_font = 'Helvetica'
            self.logger.warning("未找到中文字体，使用Helvetica")
        except Exception as e:
            self.logger.warning(f"字体设置失败: {str(e)}")
            self.default_font = 'Helvetica'
            
    def setup_styles(self):
        """设置PDF样式"""
        self.styles = getSampleStyleSheet()
        # 添加自定义样式
        self.styles.add(ParagraphStyle(
            name='Custom',
            parent=self.styles['Normal'],
            fontSize=10,
            leading=14,
            firstLineIndent=0,
            fontName=self.default_font,
            wordWrap='CJK',  # 支持中文换行
            encoding='utf-8',
            allowWidows=0,
            allowOrphans=0,
            alignment=0,  # 左对齐
        ))
        
    def clean_text(self, text):
        """清理文本内容"""
        if not text:
            return ""
            
        try:
            # 移除零宽字符和其他不可见字符
            text = re.sub(r'[\u200b-\u200f\u202a-\u202e\ufeff\u2028\u2029]', '', text)
            
            # 移除控制字符，但保留换行和制表符
            text = ''.join(char for char in text if ord(char) >= 32 or char in '\n\r\t')
            
            # 处理特殊字符
            text = text.replace('\x00', '')  # 移除 null 字符
            
            # 统一换行符
            text = text.replace('\r\n', '\n').replace('\r', '\n')
            
            # 移除连续的空行，但保留段落格式
            text = re.sub(r'\n{3,}', '\n\n', text)
            
            # 确保文本是有效的UTF-8
            text = text.encode('utf-8', errors='ignore').decode('utf-8')
            
            # 移除行尾空白字符
            lines = [line.rstrip() for line in text.split('\n')]
            text = '\n'.join(lines)
            
            return text.strip()
        except Exception as e:
            self.logger.warning(f"文本清理失败: {str(e)}")
            return str(text)
            
    def try_read_as_text(self, file_path):
        """尝试以文本方式读取文件内容"""
        try:
            # 首先尝试使用二进制模式读取文件
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                
            # 检测编码
            result = chardet.detect(raw_data)
            
            # 尝试不同的编码
            encodings = [
                result['encoding'] if result['confidence'] > 0.7 else None,
                'utf-8-sig',  # 处理带BOM的UTF-8
                'utf-8',
                'gbk',
                'gb2312',
                'gb18030',
                'big5',
                'ascii'
            ]
            
            # 过滤掉None值
            encodings = [enc for enc in encodings if enc]
            
            # 记录尝试的编码
            self.logger.info(f"尝试编码: {encodings}")
            
            for enc in encodings:
                try:
                    text = raw_data.decode(enc)
                    self.logger.info(f"成功使用编码 {enc}")
                    return self.clean_text(text)
                except Exception as e:
                    self.logger.debug(f"使用编码 {enc} 失败: {str(e)}")
                    continue
                    
            # 如果所有编码都失败，使用忽略错误的方式解码
            self.logger.warning("所有编码尝试失败，使用UTF-8（忽略错误）")
            return self.clean_text(raw_data.decode('utf-8', errors='ignore'))
            
        except Exception as e:
            raise Exception(f"无法读取文件内容: {str(e)}")
            
    def create_paragraph(self, text, style):
        """创建段落，处理特殊字符"""
        try:
            # 清理和转义文本
            text = self.clean_text(text)
            
            # 处理HTML特殊字符
            text = html.escape(text)
            
            # 处理长行
            if len(text) > 1000:  # 如果行太长，添加软换行
                text = '\n'.join(text[i:i+1000] for i in range(0, len(text), 1000))
                
            # 创建段落
            return Paragraph(text, style)
        except Exception as e:
            self.logger.warning(f"创建段落失败: {str(e)}")
            # 返回一个简单的错误提示段落
            return Paragraph("（内容无法显示）", style)
            
    def convert_unknown_file(self, input_path, output_path):
        """转换未知类型的文件"""
        try:
            # 尝试读取文件内容
            content = self.try_read_as_text(input_path)
            
            # 如果内容为空，跳过此文件
            if not content or not content.strip():
                self.logger.warning(f"文件内容为空，已跳过: {input_path}")
                return False
                
            # 创建PDF文档
            doc = SimpleDocTemplate(
                output_path,
                pagesize=A4,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72,
                encoding='utf-8'
            )
            
            # 创建内容
            story = []
            
            # 添加文件信息
            info_style = ParagraphStyle(
                'InfoStyle',
                parent=self.styles['Custom'],
                fontSize=10,
                textColor=colors.black,
                backColor=colors.lightgrey,
                borderColor=colors.black,
                borderWidth=1,
                borderPadding=5,
                spaceAfter=20,
                fontName=self.default_font
            )
            
            # 添加文件信息
            file_info = f"原始文件: {os.path.basename(input_path)}"
            story.append(Paragraph(file_info, info_style))
            story.append(Spacer(1, 20))
            
            # 添加文件内容
            content_style = ParagraphStyle(
                'ContentStyle',
                parent=self.styles['Custom'],
                fontSize=9,
                leading=12,
                fontName=self.default_font
            )
            
            # 处理内容
            valid_content = False  # 用于标记是否有有效内容
            for line in content.split('\n'):
                if line.strip():
                    try:
                        # 转义特殊字符
                        line = html.escape(line)
                        story.append(Paragraph(line, content_style))
                        story.append(Spacer(1, 6))
                        valid_content = True
                    except Exception as e:
                        self.logger.warning(f"处理行失败: {str(e)}")
                        continue
                        
            # 如果没有有效内容，返回False
            if not valid_content:
                self.logger.warning(f"文件无有效内容，已跳过: {input_path}")
                return False
                
            # 生成PDF
            doc.build(story)
            return True
        except Exception as e:
            self.logger.error(f"转换未知类型文件失败: {str(e)}")
            return False
            
    def setup_mime_types(self):
        """设置支持的MIME类型和文件扩展名"""
        self.supported_extensions = {
            # 文档
            '.txt': self.convert_text,
            '.docx': self.convert_docx,
            '.doc': self.convert_docx,
            '.xlsx': self.convert_xlsx,
            '.xls': self.convert_xlsx,
            '.pptx': self.convert_pptx,
            '.ppt': self.convert_pptx,
            # 图片
            '.jpg': self.convert_image,
            '.jpeg': self.convert_image,
            '.png': self.convert_image,
            '.gif': self.convert_image,
            '.bmp': self.convert_image,
            # 网页
            '.html': self.convert_text,
            '.htm': self.convert_text,
        }
        
    def convert_text(self, input_path, output_path):
        """转换文本文件为PDF"""
        try:
            # 读取文件内容
            content = self.try_read_as_text(input_path)
            
            # 如果内容为空，跳过此文件
            if not content or not content.strip():
                self.logger.warning(f"文件内容为空，已跳过: {input_path}")
                return False
                
            # 创建PDF文档
            doc = SimpleDocTemplate(
                output_path,
                pagesize=A4,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72,
                encoding='utf-8'
            )
            
            story = []
            style = ParagraphStyle(
                'CustomText',
                parent=self.styles['Custom'],
                fontName=self.default_font,
                fontSize=10,
                leading=14,
                spaceBefore=6,
                spaceAfter=6
            )
            
            # 处理内容
            valid_content = False  # 用于标记是否有有效内容
            for line in content.split('\n'):
                if line.strip():
                    try:
                        # 转义特殊字符
                        line = html.escape(line)
                        story.append(Paragraph(line, style))
                        valid_content = True
                    except Exception as e:
                        self.logger.warning(f"处理行失败: {str(e)}")
                        continue
                        
            # 如果没有有效内容，返回False
            if not valid_content:
                self.logger.warning(f"文件无有效内容，已跳过: {input_path}")
                return False
                
            # 生成PDF
            doc.build(story)
            return True
        except Exception as e:
            self.logger.error(f"转换文本文件失败: {str(e)}")
            return False
            
    def convert_docx(self, input_path, output_path):
        """转换DOCX文件为PDF"""
        try:
            doc = Document(input_path)
            pdf_doc = SimpleDocTemplate(
                output_path,
                pagesize=A4,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72,
                encoding='utf-8'
            )
            
            story = []
            
            # 处理段落
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    try:
                        text = html.escape(paragraph.text)
                        p = Paragraph(text, self.styles['Normal'])
                        story.append(p)
                        story.append(Spacer(1, 6))
                    except Exception as e:
                        self.logger.warning(f"处理段落失败: {str(e)}")
                        continue
                    
            # 处理表格
            for table in doc.tables:
                try:
                    data = []
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            # 转义特殊字符
                            text = html.escape(cell.text.strip())
                            row_data.append(text)
                        data.append(row_data)
                        
                    if data:
                        table_style = TableStyle([
                            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                            ('FONTNAME', (0, 0), (-1, 0), self.default_font),
                            ('FONTSIZE', (0, 0), (-1, 0), 10),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                            ('GRID', (0, 0), (-1, -1), 1, colors.black)
                        ])
                        
                        pdf_table = Table(data)
                        pdf_table.setStyle(table_style)
                        story.append(pdf_table)
                        story.append(Spacer(1, 12))
                except Exception as e:
                    self.logger.warning(f"处理表格失败: {str(e)}")
                    continue
                    
            # 生成PDF
            pdf_doc.build(story)
            return True
        except Exception as e:
            self.logger.error(f"转换DOCX文件失败: {str(e)}")
            return False
            
    def convert_xlsx(self, input_path, output_path):
        """转换XLSX文件为PDF"""
        try:
            wb = load_workbook(input_path)
            pdf_doc = SimpleDocTemplate(
                output_path,
                pagesize=A4,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72,
                encoding='utf-8'
            )
            
            story = []
            
            for sheet in wb:
                try:
                    # 添加工作表标题
                    title = Paragraph(f"<b>{html.escape(sheet.title)}</b>", self.styles['Heading1'])
                    story.append(title)
                    story.append(Spacer(1, 12))
                    
                    # 获取数据范围
                    min_row = sheet.min_row
                    max_row = min(sheet.max_row, 1000)  # 限制最大行数
                    min_col = sheet.min_column
                    max_col = min(sheet.max_column, 20)  # 限制最大列数
                    
                    # 提取数据
                    data = []
                    for row in range(min_row, max_row + 1):
                        row_data = []
                        for col in range(min_col, max_col + 1):
                            cell = sheet.cell(row, col)
                            value = str(cell.value if cell.value is not None else '')
                            # 转义特殊字符
                            value = html.escape(value)
                            row_data.append(value)
                        data.append(row_data)
                        
                    if data:
                        # 创建表格样式
                        table_style = TableStyle([
                            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                            ('FONTNAME', (0, 0), (-1, 0), self.default_font),
                            ('FONTSIZE', (0, 0), (-1, 0), 10),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('BACKGROUND', (0, 0), (-1, 0), colors.grey90),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                            ('GRID', (0, 0), (-1, -1), 1, colors.black)
                        ])
                        
                        # 创建表格
                        pdf_table = Table(data)
                        pdf_table.setStyle(table_style)
                        story.append(pdf_table)
                        story.append(Spacer(1, 20))
                except Exception as e:
                    self.logger.warning(f"处理工作表 {sheet.title} 失败: {str(e)}")
                    continue
                    
            # 生成PDF
            pdf_doc.build(story)
            return True
        except Exception as e:
            self.logger.error(f"转换XLSX文件失败: {str(e)}")
            return False
            
    def convert_pptx(self, input_path, output_path):
        """转换PPTX文件为PDF"""
        try:
            prs = Presentation(input_path)
            pdf_doc = SimpleDocTemplate(
                output_path,
                pagesize=A4,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72,
                encoding='utf-8'
            )
            
            story = []
            slide_style = ParagraphStyle(
                'SlideStyle',
                parent=self.styles['Normal'],
                fontSize=12,
                leading=14,
                spaceBefore=20,
                spaceAfter=20,
                borderWidth=1,
                borderColor=colors.black,
                borderPadding=10,
                backColor=colors.white,
                fontName=self.default_font
            )
            
            for idx, slide in enumerate(prs.slides, 1):
                try:
                    # 添加幻灯片标题
                    story.append(Paragraph(f"幻灯片 {idx}", self.styles['Heading1']))
                    story.append(Spacer(1, 12))
                    
                    # 处理形状（包括文本框）
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            # 转义特殊字符
                            text = html.escape(shape.text)
                            p = Paragraph(text, slide_style)
                            story.append(p)
                            
                    story.append(Spacer(1, 20))
                except Exception as e:
                    self.logger.warning(f"处理幻灯片 {idx} 失败: {str(e)}")
                    continue
                    
            # 生成PDF
            pdf_doc.build(story)
            return True
        except Exception as e:
            self.logger.error(f"转换PPTX文件失败: {str(e)}")
            return False
            
    def convert_image(self, input_path, output_path):
        """转换图片文件为PDF"""
        try:
            # 打开图片
            image = Image.open(input_path)
            
            # 如果是RGBA模式，转换为RGB
            if image.mode == 'RGBA':
                image = image.convert('RGB')
                
            # 计算图片在A4页面上的大小
            a4_width, a4_height = A4
            img_width, img_height = image.size
            
            # 计算缩放比例
            width_ratio = (a4_width - 2*72) / img_width  # 减去左右边距
            height_ratio = (a4_height - 2*72) / img_height  # 减去上下边距
            ratio = min(width_ratio, height_ratio)
            
            # 计算缩放后的尺寸
            new_width = img_width * ratio
            new_height = img_height * ratio
            
            # 创建PDF
            c = canvas.Canvas(output_path, pagesize=A4)
            
            # 计算居中位置
            x = (a4_width - new_width) / 2
            y = (a4_height - new_height) / 2
            
            # 将图片绘制到PDF
            c.drawImage(input_path, x, y, width=new_width, height=new_height)
            c.save()
            return True
        except Exception as e:
            self.logger.error(f"转换图片文件失败: {str(e)}")
            return False
            
    def convert_folder(self, folder_path, log_callback, progress_callback):
        """转换文件夹中的所有文件"""
        self.cancel_flag = False
        self.log_callback = log_callback
        self.progress_callback = progress_callback
        
        # 创建输出目录
        output_dir = os.path.join(folder_path, 'outputsPDF')
        os.makedirs(output_dir, exist_ok=True)
        
        # 获取所有文件
        all_files = []
        for root, _, files in os.walk(folder_path):
            # 跳过outputsPDF目录
            if 'outputsPDF' in root:
                continue
                
            for file in files:
                if not file.endswith('.pdf'):  # 排除PDF文件
                    file_path = os.path.join(root, file)
                    all_files.append(file_path)
                    
        total_files = len(all_files)
        if total_files == 0:
            self.log_callback("未找到可转换的文件！")
            return
            
        self.log_callback(f"找到 {total_files} 个文件")
        
        # 开始转换
        converted_count = 0
        for index, file_path in enumerate(all_files, 1):
            if self.cancel_flag:
                self.log_callback("转换已取消！")
                break
                
            try:
                _, ext = os.path.splitext(file_path)
                converter = self.supported_extensions.get(ext.lower())
                
                # 创建相对路径保持目录结构
                rel_path = os.path.relpath(file_path, folder_path)
                output_path = os.path.join(output_dir, rel_path + '.pdf')
                
                # 确保输出目录存在
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                
                if converter:
                    # 使用对应的转换器转换文件
                    converter(file_path, output_path)
                    converted_count += 1
                    self.log_callback(f"成功转换 ({index}/{total_files}): {rel_path}")
                else:
                    # 对于不支持的文件类型，尝试作为文本文件处理
                    self.log_callback(f"尝试读取文件内容 ({index}/{total_files}): {rel_path}")
                    if self.convert_unknown_file(file_path, output_path):
                        converted_count += 1
                        self.log_callback(f"成功转换为文本 ({index}/{total_files}): {rel_path}")
                    else:
                        self.log_callback(f"无法转换文件 ({index}/{total_files}): {rel_path}")
                    
                progress = (index / total_files) * 100
                self.progress_callback(progress)
                
            except Exception as e:
                self.log_callback(f"处理文件失败 ({index}/{total_files}): {os.path.basename(file_path)}")
                self.log_callback(f"错误信息: {str(e)}")
                
        self.log_callback(f"\n转换完成！共转换 {converted_count}/{total_files} 个文件") 